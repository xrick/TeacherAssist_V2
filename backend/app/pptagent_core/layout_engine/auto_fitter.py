import logging

from pptx.util import Pt

from app.pptagent_core.layout_engine.text_metrics import TextMetrics

logger = logging.getLogger(__name__)


class AutoFitter:
    """
    智慧型排版調整器：自動調整字體大小以適應容器。
    """

    MIN_FONT_SIZE = 8
    MAX_ITERATIONS = 12

    @classmethod
    def fit_text(
        cls,
        text_frame,
        text: str,
        font_name: str = "Arial",
        max_font_size: int = 24,
    ):
        """
        將文字填入 TextFrame，若內容過多則自動縮小字體。

        Args:
            text_frame: python-pptx 的 TextFrame 物件
            text: 文字內容
            font_name: 字體名稱
            max_font_size: 起始最大字級
        """
        if not text:
            return

        shape = text_frame._parent

        # 取得邊距 (預設值為一般 PPT 邊距)
        margin_left = getattr(text_frame, "margin_left", 91440)
        margin_right = getattr(text_frame, "margin_right", 91440)
        margin_top = getattr(text_frame, "margin_top", 45720)
        margin_bottom = getattr(text_frame, "margin_bottom", 45720)

        available_width = shape.width - margin_left - margin_right
        available_height = shape.height - margin_top - margin_bottom

        # 使用二分搜尋法 (Binary Search) 尋找最佳字級
        low = cls.MIN_FONT_SIZE
        high = max_font_size
        optimal_size = cls.MIN_FONT_SIZE

        # 先檢查最大字級是否放得下 (效能優化)
        w, h = TextMetrics.measure_text(text, font_name, high, max_width_emu=available_width)
        if h <= available_height:
            optimal_size = high
        else:
            # 開始二分搜尋
            while low <= high:
                mid = (low + high) / 2
                if mid < cls.MIN_FONT_SIZE:
                    break

                w, h = TextMetrics.measure_text(text, font_name, mid, max_width_emu=available_width)

                if h <= available_height:
                    optimal_size = mid
                    low = mid + 0.5  # 嘗試更大的
                else:
                    high = mid - 0.5  # 需要更小的

        # 填入最終內容與樣式
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(optimal_size)
        p.font.name = font_name
        text_frame.word_wrap = True

        # 若有大幅縮小，記錄 log 以便除錯
        if optimal_size < max_font_size:
            logger.debug(
                f"Auto-fit: Text shrunk to {optimal_size:.1f}pt (Space: {available_height} EMU)"
            )
