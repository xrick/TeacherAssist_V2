"""
Template Analyzer Role

Stage 1: Analyzes PPTX template structure and prepares slide skeleton

極簡原則：
- 保留設計（Masters, Layouts, Theme）
- 清除範例內容
- 建立乾淨的 placeholder 結構

v0.2 更新：
- 整合 TemplateConfig 讀取 structure_rules
- 支援 body_pool 輪替選擇 layout
"""

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import PP_PLACEHOLDER

from app.pptagent_core.config import TemplateConfig, get_template_config

logger = logging.getLogger(__name__)


class TemplateAnalyzer:
    """
    分析 PPTX Template 結構

    功能：
    1. 掃描 Template 的所有 Slide Layouts
    2. 提取 Placeholder 結構
    3. 建議適合的版面配置組合
    """

    # Placeholder 類型映射
    PH_TYPE_MAP = {
        PP_PLACEHOLDER.TITLE: "TITLE",
        PP_PLACEHOLDER.SUBTITLE: "SUBTITLE",
        PP_PLACEHOLDER.BODY: "BODY",
        PP_PLACEHOLDER.OBJECT: "CONTENT",
        PP_PLACEHOLDER.CENTER_TITLE: "CENTER_TITLE",
        PP_PLACEHOLDER.CHART: "CHART",
        PP_PLACEHOLDER.TABLE: "TABLE",
        PP_PLACEHOLDER.PICTURE: "PICTURE",
        PP_PLACEHOLDER.FOOTER: "FOOTER",
        PP_PLACEHOLDER.DATE: "DATE",
        PP_PLACEHOLDER.SLIDE_NUMBER: "SLIDE_NUMBER",
    }

    # 略過的 Placeholder 類型（自動欄位）
    SKIP_TYPES = {"FOOTER", "DATE", "SLIDE_NUMBER"}

    def __init__(self, template_path: Path, config: TemplateConfig | None = None):
        """
        初始化 Template Analyzer

        Args:
            template_path: PPTX Template 檔案路徑
            config: TemplateConfig 物件，None 則自動從 template_path 推斷
        """
        self.template_path = template_path
        self._prs: PptxPresentation | None = None
        self._layouts_cache: list[dict] | None = None

        # v0.2: 載入 config
        if config is not None:
            self.config = config
        else:
            # 從 template_path 推斷 template name
            template_name = template_path.stem  # e.g., "education_basic"
            self.config = get_template_config(template_name)

    @property
    def prs(self) -> PptxPresentation:
        """延遲載入 Presentation"""
        if self._prs is None:
            self._prs = PptxPresentation(str(self.template_path))
        return self._prs

    def get_available_layouts(self, master_index: int = 0) -> list[dict[str, Any]]:
        """
        取得所有可用的版面配置

        Args:
            master_index: 使用的 Slide Master 索引

        Returns:
            版面配置列表，每個包含 index, name, placeholders
        """
        if self._layouts_cache is not None:
            return self._layouts_cache

        master = self.prs.slide_masters[master_index]
        layouts = []

        for i, layout in enumerate(master.slide_layouts):
            placeholders = []

            for shape in layout.placeholders:
                ph_type = self._get_placeholder_type(shape.placeholder_format.type)

                # 跳過自動欄位
                if ph_type in self.SKIP_TYPES:
                    continue

                placeholders.append(
                    {
                        "idx": shape.placeholder_format.idx,
                        "type": ph_type,
                    }
                )

            layouts.append(
                {
                    "index": i,
                    "name": layout.name,
                    "placeholders": placeholders,
                }
            )

        self._layouts_cache = layouts
        return layouts

    def _get_placeholder_type(self, ph_type) -> str:
        """取得 Placeholder 類型名稱"""
        return self.PH_TYPE_MAP.get(ph_type, str(ph_type).split(".")[-1].split()[0])

    def suggest_layout_sequence(
        self,
        slide_count: int,
        include_title: bool = True,
        include_closing: bool = True,
    ) -> list[int]:
        """
        建議版面配置序列（v0.2: 使用 config 的 structure_rules）

        Args:
            slide_count: 目標投影片數量
            include_title: 是否包含標題頁
            include_closing: 是否包含結尾頁

        Returns:
            Layout 索引列表
        """
        rules = self.config.structure_rules

        # 組合序列
        sequence = []
        body_slide_count = slide_count

        # 1. Opening (標題頁)
        if include_title:
            sequence.append(rules.opening)
            body_slide_count -= 1

        # 2. Closing (結尾頁) - 先扣除數量
        if include_closing:
            body_slide_count -= 1

        # 3. Body slides - 使用 body_pool 輪替
        body_pool = rules.body_pool
        for i in range(max(0, body_slide_count)):
            layout_idx = body_pool[i % len(body_pool)]
            sequence.append(layout_idx)

        # 4. Closing (結尾頁)
        if include_closing:
            sequence.append(rules.closing)

        logger.debug(f"Layout sequence (rules): {sequence}")
        return sequence

    def _find_layout_by_type(self, layouts: list[dict], *required_types: str) -> int | None:
        """找出包含指定 placeholder 類型的 layout"""
        for layout in layouts:
            ph_types = {ph["type"] for ph in layout["placeholders"]}
            if all(t in ph_types for t in required_types):
                return layout["index"]
        return None

    def create_slide_structure(
        self,
        layout_sequence: list[int],
    ) -> dict[str, Any]:
        """
        建立投影片結構（供 LLM 填充）

        Args:
            layout_sequence: Layout 索引序列

        Returns:
            結構化的 JSON，包含每張投影片的 placeholder
        """
        layouts = self.get_available_layouts()

        structure = {
            "template": self.template_path.name,
            "slide_count": len(layout_sequence),
            "slides": [],
        }

        for slide_idx, layout_idx in enumerate(layout_sequence):
            if layout_idx >= len(layouts):
                logger.warning(f"Layout index {layout_idx} 超出範圍，使用預設")
                layout_idx = 1 if len(layouts) > 1 else 0

            layout = layouts[layout_idx]

            slide_info = {
                "index": slide_idx,
                "layout_index": layout_idx,
                "layout_name": layout["name"],
                "placeholders": [],
            }

            for ph in layout["placeholders"]:
                placeholder_info = {
                    "idx": ph["idx"],
                    "type": ph["type"],
                    "content": "",  # 待 LLM 填充
                }

                # 標記內容格式
                if ph["type"] in ("CONTENT", "BODY"):
                    placeholder_info["format"] = "bullet_list"
                    placeholder_info["content"] = []

                slide_info["placeholders"].append(placeholder_info)

            structure["slides"].append(slide_info)

        logger.info(f"建立投影片結構: {len(layout_sequence)} 張")
        return structure

    def analyze(
        self,
        slide_count: int = 10,
        include_title: bool = True,
        include_closing: bool = True,
    ) -> dict[str, Any]:
        """
        完整分析流程

        Args:
            slide_count: 目標投影片數量
            include_title: 是否包含標題頁
            include_closing: 是否包含結尾頁

        Returns:
            完整的投影片結構，可直接傳給 ContentGenerator
        """
        logger.info(f"分析 Template: {self.template_path.name}")

        # 建議版面序列
        layout_sequence = self.suggest_layout_sequence(slide_count, include_title, include_closing)

        # 建立結構
        structure = self.create_slide_structure(layout_sequence)

        logger.info(f"分析完成: {structure['slide_count']} 張投影片")
        return structure
