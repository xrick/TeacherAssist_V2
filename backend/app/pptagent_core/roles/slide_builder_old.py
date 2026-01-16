"""
Slide Builder Role

Stage 3: Builds the final PPTX file from template and content

極簡原則：
- 保留 Template 的視覺設計
- 只替換 Placeholder 內容
- 輸出乾淨的最終檔案
"""

import io
import logging
from pathlib import Path
from typing import Any

from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import PP_PLACEHOLDER

logger = logging.getLogger(__name__)


class SlideBuilder:
    """
    建構最終 PPTX 檔案

    工作流程：
    1. 載入 Template
    2. 清除現有投影片
    3. 依結構建立新投影片
    4. 填入內容
    5. 輸出 bytes
    """

    # 略過的 Placeholder 類型
    SKIP_TYPES = {
        PP_PLACEHOLDER.FOOTER,
        PP_PLACEHOLDER.DATE,
        PP_PLACEHOLDER.SLIDE_NUMBER,
    }

    def __init__(self, template_path: Path):
        """
        初始化 Slide Builder

        Args:
            template_path: PPTX Template 檔案路徑
        """
        self.template_path = template_path

    def build(
        self,
        content: dict[str, Any],
        master_index: int = 0,
    ) -> bytes:
        """
        建構 PPTX 檔案

        Args:
            content: 填入內容的結構（來自 ContentGenerator）
            master_index: 使用的 Slide Master 索引

        Returns:
            PPTX 檔案的 bytes
        """
        logger.info(f"建構 PPTX: {self.template_path.name}")

        # 載入 Template
        prs = PptxPresentation(str(self.template_path))
        master = prs.slide_masters[master_index]
        layouts = list(master.slide_layouts)

        # 清除所有現有投影片
        self._clear_slides(prs)

        # 建立新投影片並填入內容
        for slide_data in content.get("slides", []):
            layout_idx = slide_data.get("layout_index", 1)

            # 確保 layout index 有效
            if layout_idx >= len(layouts):
                logger.warning(f"Layout index {layout_idx} 超出範圍，使用預設")
                layout_idx = 1 if len(layouts) > 1 else 0

            # 新增投影片
            layout = layouts[layout_idx]
            slide = prs.slides.add_slide(layout)

            # 填入內容
            self._fill_slide(slide, slide_data)

            logger.debug(f"Slide {slide_data['index'] + 1}: {layout.name}")

        # 輸出 bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        pptx_bytes = buffer.read()

        logger.info(
            f"PPTX 建構完成: {len(content.get('slides', []))} 張投影片, {len(pptx_bytes):,} bytes"
        )

        return pptx_bytes

    def _clear_slides(self, prs: PptxPresentation) -> None:
        """清除所有現有投影片"""
        for i in range(len(prs.slides) - 1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]

        logger.debug("已清除所有現有投影片")

    def _fill_slide(self, slide, slide_data: dict[str, Any]) -> None:
        """填入單張投影片的內容"""
        # 建立 idx -> content 映射
        content_map = {}
        for ph_data in slide_data.get("placeholders", []):
            content_map[ph_data["idx"]] = ph_data

        # 遍歷所有 shape
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue

            ph_format = shape.placeholder_format

            # 跳過自動欄位
            if ph_format.type in self.SKIP_TYPES:
                continue

            ph_idx = ph_format.idx

            if ph_idx not in content_map:
                # 清除預設內容
                if shape.has_text_frame:
                    shape.text_frame.clear()
                continue

            ph_data = content_map[ph_idx]
            content_value = ph_data.get("content", "")

            if not content_value:
                if shape.has_text_frame:
                    shape.text_frame.clear()
                continue

            # 填入內容
            if shape.has_text_frame:
                self._fill_text_frame(shape.text_frame, content_value)

    def _fill_text_frame(self, text_frame, content) -> None:
        """填入 TextFrame 內容"""
        text_frame.clear()

        if isinstance(content, list):
            # 列表內容（bullet points）
            for i, item in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = str(item)
                p.level = 0
        else:
            # 單一文字
            text_frame.paragraphs[0].text = str(content)

    def build_to_file(
        self,
        content: dict[str, Any],
        output_path: Path,
        master_index: int = 0,
    ) -> Path:
        """
        建構 PPTX 檔案並存檔

        Args:
            content: 填入內容的結構
            output_path: 輸出檔案路徑
            master_index: 使用的 Slide Master 索引

        Returns:
            輸出檔案路徑
        """
        pptx_bytes = self.build(content, master_index)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_bytes(pptx_bytes)

        logger.info(f"已儲存到: {output_path}")
        return output_path
