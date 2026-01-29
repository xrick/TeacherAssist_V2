"""
Slide Builder Role

Stage 5: Builds the final PPTX file from template and content
Enhanced with Layout Engine for Auto-fit and Smart Image Placement

v0.2 更新：
- 支援 PICTURE placeholder (idx=10)
- 改進 placeholder 搜尋邏輯
"""

import io
import logging
from pathlib import Path
from typing import Any

from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches

from app.pptagent_core.config import TemplateConfig, get_template_config

# Import Layout Engine components
from app.pptagent_core.layout_engine.auto_fitter import AutoFitter

logger = logging.getLogger(__name__)


class SlideBuilder:
    """
    建構最終 PPTX 檔案

    Features:
    - Smart Layout Engine (Auto-fit text)
    - Dynamic Image Placement
    - Template Integrity Preservation
    """

    # 略過的 Placeholder 類型
    SKIP_TYPES = {
        PP_PLACEHOLDER.FOOTER,
        PP_PLACEHOLDER.DATE,
        PP_PLACEHOLDER.SLIDE_NUMBER,
    }

    def __init__(self, template_path: Path, config: TemplateConfig | None = None):
        self.template_path = template_path
        # 初始化時先設為 None，建立時再讀取
        self.slide_width = None
        self.slide_height = None

        # v0.2: 載入 config
        if config is not None:
            self.config = config
        else:
            template_name = template_path.stem
            self.config = get_template_config(template_name)

    def build(
        self,
        content: dict[str, Any],
        master_index: int = 0,
    ) -> bytes:
        logger.info(f"建構 PPTX (Enhanced): {self.template_path.name}")

        prs = PptxPresentation(str(self.template_path))

        # 讀取投影片尺寸，供排版引擎使用
        self.slide_width = prs.slide_width
        self.slide_height = prs.slide_height

        self._clear_slides(prs)

        master = prs.slide_masters[master_index]
        layouts = list(master.slide_layouts)

        # v0.3: 從 config 取得 structure_rules
        structure_rules = self.config.structure_rules if self.config else None
        body_pool = structure_rules.body_pool if structure_rules else [2]  # 預設用 TITLE_AND_BODY
        body_pool_idx = 0  # 輪替索引

        slides = content.get("slides", [])
        total_slides = len(slides)

        for i, slide_data in enumerate(slides):
            # v0.3: 智慧選擇 layout_index
            layout_idx = slide_data.get("layout_index")

            if layout_idx is None:
                # 根據 slide 位置和 layout 類型決定
                slide_layout = slide_data.get("layout", "content")

                if i == 0 or slide_layout == "title":
                    # 第一頁或明確標記為 title → 使用 opening layout
                    layout_idx = structure_rules.opening if structure_rules else 0
                elif i == total_slides - 1 or slide_layout == "closing":
                    # 最後一頁或明確標記為 closing → 使用 closing layout
                    layout_idx = structure_rules.closing if structure_rules else 0
                elif slide_layout == "section":
                    # section header
                    layout_idx = 1  # SECTION_HEADER
                else:
                    # 內容頁 → 從 body_pool 輪替
                    layout_idx = body_pool[body_pool_idx % len(body_pool)]
                    body_pool_idx += 1

            if layout_idx >= len(layouts):
                layout_idx = body_pool[0] if body_pool else 2

            layout = layouts[layout_idx]
            slide = prs.slides.add_slide(layout)

            logger.debug(f"Slide {i + 1}/{total_slides}: layout_idx={layout_idx} ({layout.name})")

            # 1. 填入文字內容 (使用 AutoFitter)
            self._fill_slide_content(slide, slide_data)

            # 2. 加入圖片 (v0.2: 傳遞 layout_index 以支援 PICTURE placeholder)
            images = slide_data.get("images", [])
            if images:
                self._place_images(
                    slide,
                    images,
                    slide_data.get("layout", "content"),
                    layout_idx,  # v0.2: 傳遞 layout_index
                )

            logger.debug(f"Slide {slide_data.get('index', '?')}: {layout.name} built")

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        pptx_bytes = buffer.read()

        logger.info(
            f"PPTX 建構完成: {len(content.get('slides', []))} 張投影片, {len(pptx_bytes):,} bytes"
        )

        return pptx_bytes

    def _clear_slides(self, prs: PptxPresentation) -> None:
        for i in range(len(prs.slides) - 1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]

    def _fill_slide_content(self, slide, slide_data: dict[str, Any]) -> None:
        """填入單張投影片的內容，使用 AutoFitter

        改進：同時支援 idx 和 type 匹配，優先使用 type 以支援不同模板
        """
        placeholders = slide_data.get("placeholders", [])

        # 建立 type -> content 映射（主要）和 idx -> content 映射（備用）
        type_map = {}
        idx_map = {}
        for ph in placeholders:
            ph_type = ph.get("type", "").upper()
            ph_idx = ph.get("idx")
            if ph_type:
                type_map[ph_type] = ph
            if ph_idx is not None:
                idx_map[ph_idx] = ph

        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue

            ph_format = shape.placeholder_format
            if ph_format.type in self.SKIP_TYPES:
                continue

            # 取得 placeholder 的 type 名稱（用於匹配）
            ph_type_name = str(ph_format.type).replace("PLACEHOLDER_TYPE.", "").replace(" (13)", "")
            ph_idx = ph_format.idx

            # 優先使用 type 匹配，其次使用 idx 匹配
            ph_data = type_map.get(ph_type_name) or type_map.get(ph_type_name.upper())

            # 特殊處理：OBJECT 類型可能對應 CONTENT 或 BODY
            if not ph_data and ph_type_name == "OBJECT":
                ph_data = type_map.get("CONTENT") or type_map.get("BODY")

            # 備用：使用 idx 匹配
            if not ph_data:
                ph_data = idx_map.get(ph_idx)

            if not ph_data:
                # 清除未使用的 placeholder 以保持整潔
                if shape.has_text_frame:
                    shape.text_frame.clear()
                continue

            content_value = ph_data.get("content", "")

            if shape.has_text_frame:
                # 將列表內容轉換為字串以便測量
                text_str = ""
                if isinstance(content_value, list):
                    text_str = "\n".join(str(x) for x in content_value)
                else:
                    text_str = str(content_value)

                # 使用 AutoFitter 替代直接賦值
                # 策略：標題固定 28pt，內文較小 (24pt)
                is_title = (
                    ph_format.type == PP_PLACEHOLDER.TITLE
                    or ph_format.type == PP_PLACEHOLDER.CENTER_TITLE
                )
                max_size = 28 if is_title else 24

                AutoFitter.fit_text(
                    shape.text_frame, text_str, font_name="Arial", max_font_size=max_size
                )

    def _place_images(self, slide, images: list, layout_type: str, layout_index: int = -1):
        """
        動態配置圖片位置

        v0.2 改進：
        - 優先使用 PICTURE placeholder (支援 idx=10)
        - 若無則使用預設位置策略

        Args:
            slide: 目標投影片
            images: 圖片資料列表
            layout_type: layout 類型字串
            layout_index: layout 索引，用於查詢 config 的 placeholder mapping
        """
        if not images:
            return

        # 暫時只處理第一張圖片，避免過度擁擠
        image_data = images[0]
        img_path = image_data.get("file_path")

        if not img_path or not Path(img_path).exists():
            logger.warning(f"Image not found: {img_path}")
            return

        # v0.2: 從 config 取得 PICTURE placeholder idx
        picture_idx = None
        if layout_index >= 0:
            ph_mapping = self.config.get_placeholder_mapping(layout_index)
            picture_idx = ph_mapping.picture  # 可能是 10 或 None
            logger.debug(f"Layout {layout_index} picture_idx from config: {picture_idx}")

        # 嘗試找到 PICTURE placeholder
        picture_placeholder = None
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_format = shape.placeholder_format
                # 優先匹配 config 指定的 idx
                if picture_idx is not None and ph_format.idx == picture_idx:
                    picture_placeholder = shape
                    break
                # 其次匹配 PICTURE type
                elif ph_format.type == PP_PLACEHOLDER.PICTURE:
                    picture_placeholder = shape
                    # 繼續搜尋，看是否有更精確的 idx 匹配

        if picture_placeholder:
            # 使用 PICTURE placeholder 的位置和尺寸
            try:
                # 取得 placeholder 的位置和尺寸
                left = picture_placeholder.left
                top = picture_placeholder.top
                width = picture_placeholder.width
                height = picture_placeholder.height

                # 使用 add_picture 放置圖片到 placeholder 的位置
                slide.shapes.add_picture(
                    str(img_path),
                    left,
                    top,
                    width=width,
                    height=height,
                )
                logger.debug(
                    f"Image placed at PICTURE placeholder position (idx={picture_placeholder.placeholder_format.idx})"
                )
                return
            except Exception as e:
                logger.warning(
                    f"Failed to use PICTURE placeholder: {e}, falling back to manual placement"
                )

        # 備用策略：手動放置圖片
        margin = Inches(0.5)

        # 根據不同的 Layout 決定圖片位置
        if layout_type == "two_column":
            # 放在右側欄位
            left = self.slide_width * 0.55
            top = Inches(2.0)
            width = self.slide_width * 0.4
            height = self.slide_height * 0.5
        elif layout_type == "image_text":
            # 圖片在左側
            left = margin
            top = Inches(1.5)
            width = self.slide_width * 0.45
            height = self.slide_height * 0.65
        elif layout_type == "full_image":
            # 全版圖片
            left = 0
            top = 0
            width = self.slide_width
            height = self.slide_height
        else:
            # 預設: 放在右下角，避免遮擋主要列表
            width = self.slide_width * 0.35
            height = self.slide_height * 0.4
            left = self.slide_width - width - margin
            top = self.slide_height - height - margin

        try:
            slide.shapes.add_picture(
                str(img_path),
                left,
                top,
                width=width,
                height=height,
            )
            logger.debug(f"Image placed manually at ({left}, {top})")
        except Exception as e:
            logger.error(f"Failed to add image: {e}")

    def build_to_file(
        self, content: dict[str, Any], output_path: Path, master_index: int = 0
    ) -> Path:
        pptx_bytes = self.build(content, master_index)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_bytes(pptx_bytes)
        return output_path
