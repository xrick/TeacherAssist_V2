#!/usr/bin/env python3
"""Debug: 檢查 PPTX 內容"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

pptx_path = Path("data/outputs/test_v03_ml.pptx")
prs = Presentation(pptx_path)

print(f"檢查 PPTX: {pptx_path}")
print(f"投影片數: {len(prs.slides)}")
print("=" * 60)

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i + 1} ---")
    print(f"Layout: {slide.slide_layout.name}")

    for shape in slide.shapes:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            ph_type = str(ph.type).replace("PLACEHOLDER_TYPE.", "")
            print(f"  Placeholder idx={ph.idx}, type={ph_type}")

            if shape.has_text_frame:
                text = shape.text_frame.text[:100] if shape.text_frame.text else "(empty)"
                print(f"    Text: {text}")
        elif hasattr(shape, "image"):
            print(f"  Image: {shape.width}x{shape.height}")
