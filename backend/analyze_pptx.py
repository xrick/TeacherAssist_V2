#!/usr/bin/env python3
"""分析 PPTX 檔案的排版結構"""

from pptx import Presentation
from pptx.util import Inches

prs = Presentation("data/outputs/test_brain.pptx")

print(f'簡報尺寸: {prs.slide_width.inches:.2f}" x {prs.slide_height.inches:.2f}"')
print(f"投影片數: {len(prs.slides)}\n")

for i, slide in enumerate(prs.slides):
    layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
    print(f"=== Slide {i + 1}: Layout '{layout_name}' ===")

    for shape in slide.shapes:
        if shape.has_text_frame:
            left = shape.left.inches if shape.left else 0
            top = shape.top.inches if shape.top else 0
            width = shape.width.inches if shape.width else 0
            height = shape.height.inches if shape.height else 0

            text = shape.text_frame.text[:100].replace("\n", " ") if shape.text_frame.text else ""

            print(f'  pos=({left:.2f}", {top:.2f}") size=({width:.2f}"x{height:.2f}")')
            print(f"  text: {text}")
            print()
