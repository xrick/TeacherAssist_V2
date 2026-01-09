#!/usr/bin/env python3
"""
rearrange.py - 從 Template 選取需要的版面配置建立工作檔

功能：
1. 讀取 template，保留所有 Masters/Layouts/Theme
2. 清除所有現有投影片
3. 依據指定的 layout 索引建立空白投影片
4. 輸出乾淨的 working.pptx

用法：
    python rearrange.py <template.pptx> <output.pptx> <layout_indices>

範例：
    python rearrange.py creative_colorful.pptx working.pptx "1,2,2,2,2,2,2,2,2,5"
    # 建立 10 張投影片：1張標題頁 + 8張內容頁 + 1張結尾頁
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def list_layouts(template_path: str) -> None:
    """列出 template 中所有可用的版面配置"""
    prs = Presentation(template_path)

    print(f"\n📋 Template: {template_path}")
    print("=" * 60)

    for master_idx, master in enumerate(prs.slide_masters):
        master_name = master.name or f"Master_{master_idx}"
        print(f"\n🎨 Master {master_idx}: {master_name}")
        print("-" * 40)

        for i, layout in enumerate(master.slide_layouts):
            # 收集 placeholder 類型
            phs = []
            for shape in layout.placeholders:
                ph_type = str(shape.placeholder_format.type)
                ph_type = ph_type.split(".")[-1].split()[0]
                phs.append(ph_type)

            ph_str = ", ".join(phs) if phs else "(blank)"
            print(f"  [{i:2d}] {layout.name:<30} | {ph_str}")


def rearrange(
    template_path: str, output_path: str, layout_indices: list[int], master_index: int = 0
) -> None:
    """
    從 template 建立新的 PPTX，只包含指定的版面配置

    Args:
        template_path: 來源 template 檔案路徑
        output_path: 輸出檔案路徑
        layout_indices: 要使用的 layout 索引列表（決定投影片數量和版面）
        master_index: 使用的 master 索引（預設 0）
    """
    # 載入 template
    prs = Presentation(template_path)

    # 取得指定的 master
    if master_index >= len(prs.slide_masters):
        raise ValueError(f"Master index {master_index} 超出範圍 (共 {len(prs.slide_masters)} 個)")

    master = prs.slide_masters[master_index]
    layouts = list(master.slide_layouts)

    # 驗證 layout indices
    for idx in layout_indices:
        if idx >= len(layouts):
            raise ValueError(f"Layout index {idx} 超出範圍 (共 {len(layouts)} 個)")

    # 刪除所有現有投影片（從後往前刪）
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    # 依據 layout_indices 建立新投影片
    for i, layout_idx in enumerate(layout_indices):
        layout = layouts[layout_idx]
        slide = prs.slides.add_slide(layout)

        # 清除所有 placeholder 中的預設文字
        for shape in slide.shapes:
            if shape.is_placeholder and shape.has_text_frame:
                shape.text_frame.clear()

        print(f"  ✓ Slide {i + 1}: Layout[{layout_idx}] = {layout.name}")

    # 儲存
    prs.save(output_path)
    print(f"\n✅ 已建立 {output_path} ({len(layout_indices)} 張投影片)")


def main():
    parser = argparse.ArgumentParser(
        description="從 PPTX Template 選取版面配置建立工作檔",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
範例：
  # 列出所有版面配置
  python rearrange.py template.pptx --list

  # 建立 10 張投影片（標題 + 8內容 + 結尾）
  python rearrange.py template.pptx output.pptx "1,2,2,2,2,2,2,2,2,5"

  # 建立簡單的 5 張投影片
  python rearrange.py template.pptx output.pptx "1,2,2,2,1"
        """,
    )

    parser.add_argument("template", help="來源 template PPTX 檔案")
    parser.add_argument("output", nargs="?", help="輸出 PPTX 檔案路徑")
    parser.add_argument("layouts", nargs="?", help="Layout 索引，以逗號分隔 (例: 1,2,2,2,5)")
    parser.add_argument("--list", "-l", action="store_true", help="列出所有可用的版面配置")
    parser.add_argument("--master", "-m", type=int, default=0, help="使用的 Master 索引 (預設: 0)")

    args = parser.parse_args()

    # 檢查 template 存在
    if not Path(args.template).exists():
        print(f"❌ 找不到檔案: {args.template}")
        sys.exit(1)

    # 列出模式
    if args.list:
        list_layouts(args.template)
        sys.exit(0)

    # 建立模式 - 需要 output 和 layouts
    if not args.output or not args.layouts:
        parser.print_help()
        print("\n❌ 需要指定 output 和 layouts 參數")
        sys.exit(1)

    # 解析 layout indices
    try:
        layout_indices = [int(x.strip()) for x in args.layouts.split(",")]
    except ValueError:
        print(f"❌ 無效的 layouts 格式: {args.layouts}")
        print("   請使用逗號分隔的數字，例如: 1,2,2,2,5")
        sys.exit(1)

    print(f"\n🔧 Rearrange PPTX")
    print(f"   Template: {args.template}")
    print(f"   Output: {args.output}")
    print(f"   Layouts: {layout_indices}")
    print()

    rearrange(args.template, args.output, layout_indices, args.master)


if __name__ == "__main__":
    main()
