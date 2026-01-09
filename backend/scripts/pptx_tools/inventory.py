#!/usr/bin/env python3
"""
inventory.py - 提取 PPTX 結構，生成 JSON 清單供 LLM 填充

功能：
1. 讀取 working.pptx（由 rearrange.py 產生）
2. 提取每張投影片的 placeholder 結構
3. 輸出 JSON 格式的「填空模板」

用法：
    python inventory.py <input.pptx> [output.json]

輸出格式：
{
  "template": "creative_colorful.pptx",
  "slide_count": 10,
  "slides": [
    {
      "index": 0,
      "layout": "Title Slide",
      "placeholders": [
        {"id": 0, "type": "TITLE", "content": ""},
        {"id": 1, "type": "SUBTITLE", "content": ""}
      ]
    },
    ...
  ]
}
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


def get_placeholder_type_name(ph_type) -> str:
    """取得 placeholder 類型的可讀名稱"""
    type_map = {
        PP_PLACEHOLDER.TITLE: "TITLE",
        PP_PLACEHOLDER.SUBTITLE: "SUBTITLE",
        PP_PLACEHOLDER.BODY: "BODY",
        PP_PLACEHOLDER.OBJECT: "CONTENT",
        PP_PLACEHOLDER.CENTER_TITLE: "CENTER_TITLE",
        PP_PLACEHOLDER.CHART: "CHART",
        PP_PLACEHOLDER.TABLE: "TABLE",
        PP_PLACEHOLDER.PICTURE: "PICTURE",
        PP_PLACEHOLDER.FOOTER: "FOOTER",
        PP_PLACEHOLDER.DATE: "DATE",
        PP_PLACEHOLDER.SLIDE_NUMBER: "SLIDE_NUMBER",
    }
    return type_map.get(ph_type, str(ph_type).split(".")[-1].split()[0])


def extract_inventory(pptx_path: str) -> dict:
    """
    提取 PPTX 的 placeholder 結構

    Args:
        pptx_path: PPTX 檔案路徑

    Returns:
        結構化的 JSON 資料
    """
    prs = Presentation(pptx_path)

    inventory = {"source_file": Path(pptx_path).name, "slide_count": len(prs.slides), "slides": []}

    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {"index": slide_idx, "layout": slide.slide_layout.name, "placeholders": []}

        # 收集所有可編輯的 placeholder
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue

            ph_format = shape.placeholder_format
            ph_type = get_placeholder_type_name(ph_format.type)

            # 跳過頁碼、日期、頁尾等自動欄位
            if ph_type in ("SLIDE_NUMBER", "DATE", "FOOTER"):
                continue

            placeholder_info = {
                "idx": ph_format.idx,
                "type": ph_type,
                "content": "",  # 空白，待 LLM 填充
            }

            # 如果有多行內容的 placeholder，標記為 list 類型
            if ph_type in ("CONTENT", "BODY"):
                placeholder_info["format"] = "bullet_list"
                placeholder_info["content"] = []  # 用列表格式

            slide_info["placeholders"].append(placeholder_info)

        # 依 idx 排序
        slide_info["placeholders"].sort(key=lambda x: x["idx"])

        inventory["slides"].append(slide_info)

    return inventory


def print_inventory(inventory: dict) -> None:
    """以易讀格式印出結構"""
    print(f"\n📊 PPTX 結構清單")
    print(f"   來源: {inventory['source_file']}")
    print(f"   投影片數: {inventory['slide_count']}")
    print("=" * 50)

    for slide in inventory["slides"]:
        print(f"\n📑 Slide {slide['index'] + 1}: {slide['layout']}")
        for ph in slide["placeholders"]:
            format_hint = f" [{ph.get('format', 'text')}]" if ph.get("format") else ""
            print(f"   • [{ph['idx']}] {ph['type']}{format_hint}")


def main():
    parser = argparse.ArgumentParser(
        description="提取 PPTX 結構，生成 JSON 清單",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
範例：
  # 印出結構
  python inventory.py working.pptx

  # 輸出 JSON 檔案
  python inventory.py working.pptx structure.json

  # 輸出到 stdout（供 pipe 使用）
  python inventory.py working.pptx --stdout
        """,
    )

    parser.add_argument("input", help="輸入 PPTX 檔案")
    parser.add_argument("output", nargs="?", help="輸出 JSON 檔案路徑")
    parser.add_argument("--stdout", "-s", action="store_true", help="輸出 JSON 到 stdout")
    parser.add_argument("--pretty", "-p", action="store_true", help="格式化 JSON 輸出")

    args = parser.parse_args()

    # 檢查輸入檔案
    if not Path(args.input).exists():
        print(f"❌ 找不到檔案: {args.input}")
        sys.exit(1)

    # 提取結構
    inventory = extract_inventory(args.input)

    # 輸出模式
    if args.stdout:
        # 輸出到 stdout
        indent = 2 if args.pretty else None
        print(json.dumps(inventory, ensure_ascii=False, indent=indent))
    elif args.output:
        # 輸出到檔案
        with open(args.output, "w", encoding="utf-8") as f:
            json.dump(inventory, f, ensure_ascii=False, indent=2)
        print(f"✅ 已輸出到 {args.output}")
        print_inventory(inventory)
    else:
        # 印出易讀格式
        print_inventory(inventory)
        print("\n💡 提示: 使用 --stdout 或指定輸出檔案以取得 JSON 格式")


if __name__ == "__main__":
    main()
