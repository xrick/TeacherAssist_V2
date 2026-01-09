#!/usr/bin/env python3
"""
replace.py - 將 LLM 生成的內容填入 PPTX

功能：
1. 讀取 working.pptx（由 rearrange.py 產生）
2. 讀取 content.json（由 LLM 生成，或使用 --generate 自動生成）
3. 將內容填入對應的 placeholder
4. 輸出最終 output.pptx

用法：
    # 使用現有 JSON
    python replace.py working.pptx content.json output.pptx

    # 自動生成內容（需設定 OPENAI_API_KEY）
    python replace.py working.pptx output.pptx --generate "機器學習入門"

JSON 內容格式（與 inventory.py 輸出相同，但填入 content）：
{
  "slides": [
    {
      "index": 0,
      "placeholders": [
        {"idx": 0, "type": "TITLE", "content": "機器學習入門"},
        {"idx": 1, "type": "SUBTITLE", "content": "從零開始理解 AI"}
      ]
    },
    ...
  ]
}
"""

import argparse
import json
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

# LLM 生成的 System Prompt
SYSTEM_PROMPT = """你是一位專業的簡報內容生成專家。請根據提供的主題和投影片結構，生成簡潔、專業的內容。

原則：
1. 簡約清晰：每個要點控制在 15 字以內
2. 一目瞭然：使用具體數據或例子，避免抽象描述
3. 結構分明：每張投影片只講一個重點
4. 專業術語：保留英文專業術語，不要強行翻譯

格式要求：
- TITLE: 簡短標題（5-10 字）
- SUBTITLE: 副標題或說明（10-20 字）
- CONTENT/BODY: 3-5 個要點，每點一行

請直接返回 JSON，不要加任何說明。"""


def generate_content_with_llm(topic: str, inventory: dict) -> dict:
    """
    使用 LLM 生成簡報內容

    Args:
        topic: 簡報主題
        inventory: 從 inventory.py 提取的結構

    Returns:
        填入內容的 inventory
    """
    try:
        from openai import OpenAI
    except ImportError:
        print("❌ 需要安裝 openai: pip install openai")
        sys.exit(1)

    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        print("❌ 請設定 OPENAI_API_KEY 環境變數")
        sys.exit(1)

    client = OpenAI(api_key=api_key)

    # 建立 prompt
    user_prompt = f"""主題：{topic}

請為以下投影片結構生成內容：

{json.dumps(inventory, ensure_ascii=False, indent=2)}

請填入每個 placeholder 的 content 欄位，返回完整 JSON。
對於 bullet_list 格式的 content，請返回字串列表。"""

    print("🤖 正在使用 LLM 生成內容...")

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.7,
        max_tokens=4000,
    )

    content = response.choices[0].message.content

    # 解析 JSON
    try:
        # 嘗試直接解析
        result = json.loads(content)
    except json.JSONDecodeError:
        # 嘗試從 code block 中提取
        if "```json" in content:
            start = content.find("```json") + 7
            end = content.find("```", start)
            result = json.loads(content[start:end].strip())
        elif "```" in content:
            start = content.find("```") + 3
            end = content.find("```", start)
            result = json.loads(content[start:end].strip())
        else:
            # 找 JSON 物件
            start = content.find("{")
            end = content.rfind("}") + 1
            result = json.loads(content[start:end])

    print(f"✅ 內容生成完成")
    return result


def apply_content(pptx_path: str, content: dict, output_path: str) -> None:
    """
    將內容填入 PPTX

    Args:
        pptx_path: 來源 PPTX 檔案
        content: 包含內容的 JSON 結構
        output_path: 輸出檔案路徑
    """
    prs = Presentation(pptx_path)

    slides_content = content.get("slides", [])

    for slide_data in slides_content:
        slide_idx = slide_data["index"]

        if slide_idx >= len(prs.slides):
            print(f"⚠️ 跳過不存在的投影片 {slide_idx}")
            continue

        slide = prs.slides[slide_idx]

        # 建立 idx -> content 的映射
        content_map = {}
        for ph_data in slide_data.get("placeholders", []):
            content_map[ph_data["idx"]] = ph_data

        # 填入內容
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue

            ph_idx = shape.placeholder_format.idx

            if ph_idx not in content_map:
                continue

            ph_data = content_map[ph_idx]
            content_value = ph_data.get("content", "")

            if not content_value:
                continue

            # 處理不同類型的內容
            if shape.has_text_frame:
                tf = shape.text_frame
                tf.clear()

                if isinstance(content_value, list):
                    # 列表內容（bullet points）
                    for i, item in enumerate(content_value):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = str(item)
                        p.level = 0
                else:
                    # 單一文字
                    tf.paragraphs[0].text = str(content_value)

        print(f"  ✓ Slide {slide_idx + 1}: {slide.slide_layout.name}")

    # 儲存
    prs.save(output_path)
    print(f"\n✅ 已儲存到 {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description="將 LLM 生成的內容填入 PPTX",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
範例：
  # 使用現有 JSON 檔案
  python replace.py working.pptx content.json output.pptx

  # 自動生成內容
  python replace.py working.pptx output.pptx --generate "機器學習入門"

  # 生成並儲存 JSON（不套用到 PPTX）
  python replace.py working.pptx --generate "機器學習" --save-json content.json
        """,
    )

    parser.add_argument("input", help="輸入 PPTX 檔案（working.pptx）")
    parser.add_argument(
        "content_or_output", help="JSON 內容檔案，或輸出 PPTX 路徑（配合 --generate）"
    )
    parser.add_argument("output", nargs="?", help="輸出 PPTX 檔案路徑")
    parser.add_argument("--generate", "-g", metavar="TOPIC", help="使用 LLM 自動生成指定主題的內容")
    parser.add_argument("--save-json", "-j", metavar="FILE", help="儲存生成的 JSON 內容")

    args = parser.parse_args()

    # 檢查輸入檔案
    if not Path(args.input).exists():
        print(f"❌ 找不到檔案: {args.input}")
        sys.exit(1)

    # 確定模式
    if args.generate:
        # 自動生成模式
        output_path = args.content_or_output

        # 先提取結構
        from inventory import extract_inventory

        inventory = extract_inventory(args.input)

        # 生成內容
        content = generate_content_with_llm(args.generate, inventory)

        # 儲存 JSON（如果指定）
        if args.save_json:
            with open(args.save_json, "w", encoding="utf-8") as f:
                json.dump(content, f, ensure_ascii=False, indent=2)
            print(f"💾 JSON 已儲存到 {args.save_json}")

        # 套用內容
        print(f"\n🔧 套用內容到 PPTX")
        apply_content(args.input, content, output_path)

    else:
        # 使用現有 JSON 模式
        content_path = args.content_or_output
        output_path = args.output

        if not output_path:
            parser.print_help()
            print("\n❌ 請指定輸出檔案路徑")
            sys.exit(1)

        if not Path(content_path).exists():
            print(f"❌ 找不到內容檔案: {content_path}")
            sys.exit(1)

        # 讀取 JSON
        with open(content_path, "r", encoding="utf-8") as f:
            content = json.load(f)

        print(f"\n🔧 套用內容到 PPTX")
        print(f"   來源: {args.input}")
        print(f"   內容: {content_path}")
        print(f"   輸出: {output_path}")
        print()

        apply_content(args.input, content, output_path)


if __name__ == "__main__":
    main()
