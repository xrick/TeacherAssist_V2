import copy

from pptx import Presentation


def create_presentation_from_clean_template(output_filename):
    # 1. 載入您的清潔版範本
    # 注意：這裡假設您上傳的檔案名為 'professional_corporate.pptx'
    prs = Presentation("professional_corporate.pptx")

    # 定義版型索引 (對應您的檔案結構)
    LAYOUT_COVER = 0  # Slide 1
    LAYOUT_AGENDA = 1  # Slide 2
    LAYOUT_CONTENT = 2  # Slide 3
    LAYOUT_ENDING = 3  # Slide 4

    # 清空範本中預留的投影片，只保留 Master Layouts
    # (實務上 python-pptx 是使用 slide_layouts 來新增，所以不需要刪除舊 slide，
    # 但我們要確保新增時是用對應的 layout)

    # ------------------------------------------
    # 2. 開始生成新簡報
    # ------------------------------------------

    # A. 建立封面 (使用 Layout 0)
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_COVER])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]  # 通常副標題是 index 1

    title.text = "2025 年度專案計畫"
    subtitle.text = "自動化生成系統架構報告"

    # B. 建立大綱 (使用 Layout 1)
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_AGENDA])
    title = slide.shapes.title
    title.text = "今日議程"
    # 這裡通常需要更複雜的邏輯來處理條列式清單，這裡簡化演示
    content = slide.placeholders[1]
    content.text = "1. 專案背景\n2. 技術架構\n3. 時程規劃"

    # C. 建立內容頁 (使用 Layout 2) - 這是會被重複使用最多次的
    # 模擬生成 2 頁內容
    for i in range(1, 3):
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])

        # 設定標題
        title = slide.shapes.title
        title.text = f"章節 {i}: 詳細技術說明"

        # 設定內文 (對應 Lorem ipsum )
        body = slide.placeholders[1]
        body.text = "這裡是自動生成的詳細內容。\n我們保留了原始範本的字體大小與顏色配置。\n\n- 重點一\n- 重點二"

        # 若有圖片，會抓取 Image Placeholder 進行替換 (需額外處理)

    # D. 建立結尾 (使用 Layout 3)
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_ENDING])
    title = slide.shapes.title
    title.text = "報告結束，謝謝大家！"

    # 3. 存檔
    prs.save(output_filename)
    print(f"成功生成: {output_filename}")


# 執行範例
# create_presentation_from_clean_template("output_presentation.pptx")
