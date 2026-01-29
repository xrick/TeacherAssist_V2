#!/usr/bin/env python3
"""
rebuild_template.py

從 education_basic.pptx 建立修復版 education_minimal.pptx：
1. 修復 placeholder idx 衝突（所有 idx=0 → 唯一值）
2. 新增 Full Image layout（全版圖片）
3. 保持原有 color theme

執行方式：
    cd backend/data/templates
    python rebuild_template.py
"""

import shutil
import tempfile
import zipfile
from pathlib import Path

from lxml import etree

# 路徑設定
SCRIPT_DIR = Path(__file__).parent
SOURCE_PPTX = SCRIPT_DIR / "education_basic.pptx"
OUTPUT_PPTX = SCRIPT_DIR / "education_minimal.pptx"

# XML namespaces
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

# EMU 常數
EMU_PER_INCH = 914400
SLIDE_WIDTH = int(13.333 * EMU_PER_INCH)
SLIDE_HEIGHT = int(7.5 * EMU_PER_INCH)


def fix_placeholder_idx(layout_xml_bytes: bytes, layout_name: str) -> bytes:
    """
    修復單一 layout XML 中的 placeholder idx 衝突

    策略：
    - TITLE/CENTER_TITLE → idx=0
    - SUBTITLE → idx=1
    - OBJECT/BODY 依序 → idx=10, 11, 12...
    - FOOTER/DATE/SLIDE_NUMBER 保持原值
    """
    root = etree.fromstring(layout_xml_bytes)

    object_counter = 10

    for sp in root.findall(".//p:sp", NS):
        nvSpPr = sp.find("p:nvSpPr", NS)
        if nvSpPr is None:
            continue
        nvPr = nvSpPr.find("p:nvPr", NS)
        if nvPr is None:
            continue
        ph = nvPr.find("p:ph", NS)
        if ph is None:
            continue

        ph_type = ph.get("type", "")

        # 根據 type 分配 idx
        if ph_type in ("title", "ctrTitle"):
            ph.set("idx", "0")
        elif ph_type == "subTitle":
            ph.set("idx", "1")
        elif ph_type in ("body", "") or ph_type is None:
            # OBJECT type 在 XML 中可能沒有 type 屬性
            ph.set("idx", str(object_counter))
            object_counter += 1
        # FOOTER(15), DATE(16), SLIDE_NUMBER(13) 保持原值

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def create_full_image_layout() -> bytes:
    """
    建立 Full Image layout XML
    只有一個全版 PICTURE placeholder
    """
    xml_str = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
             xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
             xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
             type="blank" preserve="1">
  <p:cSld name="Full Image">
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="0" y="0"/>
          <a:ext cx="0" cy="0"/>
          <a:chOff x="0" y="0"/>
          <a:chExt cx="0" cy="0"/>
        </a:xfrm>
      </p:grpSpPr>
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="2" name="Picture Placeholder 1"/>
          <p:cNvSpPr>
            <a:spLocks noGrp="1"/>
          </p:cNvSpPr>
          <p:nvPr>
            <p:ph type="pic" idx="1"/>
          </p:nvPr>
        </p:nvSpPr>
        <p:spPr>
          <a:xfrm>
            <a:off x="0" y="0"/>
            <a:ext cx="{SLIDE_WIDTH}" cy="{SLIDE_HEIGHT}"/>
          </a:xfrm>
          <a:prstGeom prst="rect">
            <a:avLst/>
          </a:prstGeom>
        </p:spPr>
        <p:txBody>
          <a:bodyPr/>
          <a:lstStyle/>
          <a:p>
            <a:endParaRPr lang="en-US"/>
          </a:p>
        </p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr>
    <a:masterClrMapping/>
  </p:clrMapOvr>
</p:sldLayout>'''
    return xml_str.encode("utf-8")


def create_layout_rels() -> bytes:
    """建立 layout 的 rels 檔案"""
    xml_str = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/>
</Relationships>"""
    return xml_str.encode("utf-8")


def rebuild_template():
    """主函數：重建模板"""
    print("=" * 60)
    print("Rebuilding Template")
    print("=" * 60)

    if not SOURCE_PPTX.exists():
        print(f"ERROR: Source not found: {SOURCE_PPTX}")
        return False

    print(f"Source: {SOURCE_PPTX}")
    print(f"Output: {OUTPUT_PPTX}")

    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir = Path(tmpdir)

        # 解壓縮
        print("\n1. Extracting PPTX...")
        with zipfile.ZipFile(SOURCE_PPTX, "r") as zf:
            zf.extractall(tmpdir)

        # 修復 layout idx
        print("\n2. Fixing placeholder idx...")
        layouts_dir = tmpdir / "ppt" / "slideLayouts"
        for layout_file in sorted(layouts_dir.glob("slideLayout*.xml")):
            if "_rels" in str(layout_file):
                continue

            with open(layout_file, "rb") as f:
                original = f.read()

            fixed = fix_placeholder_idx(original, layout_file.name)

            with open(layout_file, "wb") as f:
                f.write(fixed)

            print(f"   Fixed: {layout_file.name}")

        # 新增 Full Image layout
        print("\n3. Adding Full Image layout...")
        new_layout_num = 13
        new_layout_path = layouts_dir / f"slideLayout{new_layout_num}.xml"
        new_rels_dir = layouts_dir / "_rels"
        new_rels_path = new_rels_dir / f"slideLayout{new_layout_num}.xml.rels"

        with open(new_layout_path, "wb") as f:
            f.write(create_full_image_layout())
        print(f"   Created: slideLayout{new_layout_num}.xml")

        new_rels_dir.mkdir(exist_ok=True)
        with open(new_rels_path, "wb") as f:
            f.write(create_layout_rels())
        print(f"   Created: slideLayout{new_layout_num}.xml.rels")

        # 更新 slideMaster rels
        print("\n4. Updating slideMaster rels...")
        master_rels_path = tmpdir / "ppt" / "slideMasters" / "_rels" / "slideMaster1.xml.rels"
        with open(master_rels_path, "rb") as f:
            rels_content = f.read()

        rels_root = etree.fromstring(rels_content)

        # 找最大 rId
        max_rid = 0
        for rel in rels_root:
            rid = rel.get("Id", "rId0")
            if rid.startswith("rId"):
                num = int(rid[3:])
                max_rid = max(max_rid, num)

        new_rid = f"rId{max_rid + 1}"

        # 新增 relationship
        new_rel = etree.SubElement(rels_root, "Relationship")
        new_rel.set("Id", new_rid)
        new_rel.set(
            "Type",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
        )
        new_rel.set("Target", f"../slideLayouts/slideLayout{new_layout_num}.xml")

        with open(master_rels_path, "wb") as f:
            f.write(
                etree.tostring(rels_root, xml_declaration=True, encoding="UTF-8", standalone=True)
            )
        print(f"   Added: {new_rid} -> slideLayout{new_layout_num}.xml")

        # 更新 slideMaster sldLayoutIdLst
        print("\n5. Updating slideMaster layout list...")
        master_path = tmpdir / "ppt" / "slideMasters" / "slideMaster1.xml"
        with open(master_path, "rb") as f:
            master_content = f.read()

        master_root = etree.fromstring(master_content)

        sldLayoutIdLst = master_root.find(".//p:sldLayoutIdLst", NS)
        if sldLayoutIdLst is not None:
            # 找最大 id
            max_id = 0
            for elem in sldLayoutIdLst:
                id_val = elem.get("id", "0")
                max_id = max(max_id, int(id_val))

            # 新增 sldLayoutId
            new_layout_id = etree.SubElement(sldLayoutIdLst, f"{{{NS['p']}}}sldLayoutId")
            new_layout_id.set("id", str(max_id + 1))
            new_layout_id.set(f"{{{NS['r']}}}id", new_rid)

            with open(master_path, "wb") as f:
                f.write(
                    etree.tostring(
                        master_root, xml_declaration=True, encoding="UTF-8", standalone=True
                    )
                )
            print(f"   Added: sldLayoutId id={max_id + 1}")

        # 更新 [Content_Types].xml
        print("\n6. Updating Content_Types...")
        ct_path = tmpdir / "[Content_Types].xml"
        with open(ct_path, "rb") as f:
            ct_content = f.read()

        ct_root = etree.fromstring(ct_content)
        ct_ns = ct_root.nsmap.get(
            None, "http://schemas.openxmlformats.org/package/2006/content-types"
        )

        new_override = etree.SubElement(ct_root, f"{{{ct_ns}}}Override")
        new_override.set("PartName", f"/ppt/slideLayouts/slideLayout{new_layout_num}.xml")
        new_override.set(
            "ContentType",
            "application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml",
        )

        with open(ct_path, "wb") as f:
            f.write(
                etree.tostring(ct_root, xml_declaration=True, encoding="UTF-8", standalone=True)
            )
        print(f"   Added: slideLayout{new_layout_num}.xml")

        # 重新打包
        print("\n7. Repacking PPTX...")
        with zipfile.ZipFile(OUTPUT_PPTX, "w", zipfile.ZIP_DEFLATED) as zf:
            for file in tmpdir.rglob("*"):
                if file.is_file():
                    arcname = file.relative_to(tmpdir)
                    zf.write(file, arcname)

        print(f"\n{'=' * 60}")
        print(f"SUCCESS: {OUTPUT_PPTX}")
        print(f"Size: {OUTPUT_PPTX.stat().st_size:,} bytes")
        print(f"{'=' * 60}")

    return True


def verify_output():
    """驗證輸出"""
    print("\nVerifying output...")

    from pptx import Presentation

    prs = Presentation(str(OUTPUT_PPTX))

    print(f"Layouts: {len(prs.slide_layouts)}")

    # 檢查 idx 唯一性
    issues = 0
    for i, layout in enumerate(prs.slide_layouts):
        idx_set = set()
        for ph in layout.placeholders:
            idx = ph.placeholder_format.idx
            ph_type = str(ph.placeholder_format.type).split(".")[-1]
            if idx in idx_set and ph_type not in ("FOOTER", "DATE", "SLIDE_NUMBER"):
                issues += 1
            idx_set.add(idx)

    if issues == 0:
        print("OK: No idx conflicts")
    else:
        print(f"WARNING: {issues} idx conflicts found")

    # 顯示前幾個 layout
    print("\nSample layouts:")
    for i in [0, 1, 2, 12]:
        if i < len(prs.slide_layouts):
            layout = prs.slide_layouts[i]
            phs = [
                (ph.placeholder_format.idx, str(ph.placeholder_format.type).split(".")[-1])
                for ph in layout.placeholders
                if str(ph.placeholder_format.type).split(".")[-1]
                not in ("FOOTER (15)", "DATE (16)", "SLIDE_NUMBER (13)")
            ]
            print(f"  {i}: {layout.name} -> {phs}")


if __name__ == "__main__":
    if rebuild_template():
        verify_output()
