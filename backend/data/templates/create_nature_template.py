"""
Create Nature Artistic style PPTX template
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

OUTPUT_PATH = Path(__file__).parent / "nature_artistic.pptx"


def create_nature_template():
    """Create a nature-themed PPTX template"""
    prs = Presentation()

    # Set slide size (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Get the blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # Create a sample slide
    slide = prs.slides.add_slide(blank_layout)

    # Add a green header shape
    header = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(13.333), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(22, 101, 52)  # #166534
    header.line.fill.background()

    # Add title text box
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(10), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "自然藝術風格"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Add decorative bottom shape (simulating nature/hills)
    bottom = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(6.5),
        Inches(13.333), Inches(1)
    )
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = RGBColor(134, 239, 172)  # #86efac
    bottom.line.fill.background()

    # Save the presentation
    prs.save(OUTPUT_PATH)
    print(f"Created: {OUTPUT_PATH}")


if __name__ == "__main__":
    create_nature_template()
