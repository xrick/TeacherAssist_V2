"""
Stage 5 Integration Test: SlideBuilder

驗證 SlideBuilder.build() 正確建構 PPTX 檔案。
"""

import io
from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from app.pptagent_core.roles.slide_builder import SlideBuilder


@pytest.fixture
def template_path() -> Path:
    """取得測試用 Template 路徑"""
    path = Path(__file__).parent.parent.parent / "data" / "templates" / "modern_clean.pptx"
    if not path.exists():
        pytest.skip(f"Template 不存在: {path}")
    return path


@pytest.fixture
def sample_enriched_content():
    """模擬 Stage 4 輸出 - placeholders 為 list of dicts"""
    return {
        "slides": [
            {
                "index": 0,
                "layout_index": 0,
                "layout_name": "Title Slide",
                "placeholders": [
                    {"idx": 0, "type": "CENTER_TITLE", "content": "機器學習入門"},
                    {"idx": 1, "type": "SUBTITLE", "content": "人工智慧基礎"},
                ],
                "images": [],
            },
            {
                "index": 1,
                "layout_index": 1,
                "layout_name": "Title and Content",
                "placeholders": [
                    {"idx": 0, "type": "TITLE", "content": "什麼是機器學習？"},
                    {
                        "idx": 1,
                        "type": "BODY",
                        "content": "• 從資料中學習\n• 自動改進效能\n• 無需明確程式設計",
                    },
                ],
                "images": [],
            },
            {
                "index": 2,
                "layout_index": 1,
                "layout_name": "Title and Content",
                "placeholders": [
                    {"idx": 0, "type": "TITLE", "content": "機器學習類型"},
                    {
                        "idx": 1,
                        "type": "BODY",
                        "content": "• 監督式學習\n• 非監督式學習\n• 強化學習",
                    },
                ],
                "images": [],
            },
        ]
    }


class TestStage5SlideBuilder:
    """Stage 5: SlideBuilder 輸出驗證"""

    def test_build_returns_bytes(self, template_path: Path, sample_enriched_content):
        """測試 build() 回傳 bytes"""
        builder = SlideBuilder(template_path)
        result = builder.build(sample_enriched_content)

        assert isinstance(result, bytes), "結果應為 bytes"
        assert len(result) > 0, "PPTX bytes 不應為空"

    def test_pptx_is_valid(self, template_path: Path, sample_enriched_content):
        """測試 PPTX 可被解析"""
        builder = SlideBuilder(template_path)
        pptx_bytes = builder.build(sample_enriched_content)

        try:
            prs = PptxPresentation(io.BytesIO(pptx_bytes))
            assert prs is not None
        except Exception as e:
            pytest.fail(f"無法解析 PPTX: {e}")

    def test_pptx_has_correct_slide_count(self, template_path: Path, sample_enriched_content):
        """測試 PPTX 有正確的投影片數量"""
        builder = SlideBuilder(template_path)
        pptx_bytes = builder.build(sample_enriched_content)

        prs = PptxPresentation(io.BytesIO(pptx_bytes))
        expected = len(sample_enriched_content["slides"])
        actual = len(prs.slides)

        assert actual == expected, f"投影片數量應為 {expected}，實際: {actual}"

    def test_pptx_slides_have_text(self, template_path: Path, sample_enriched_content):
        """測試 PPTX 投影片有文字內容"""
        builder = SlideBuilder(template_path)
        pptx_bytes = builder.build(sample_enriched_content)

        prs = PptxPresentation(io.BytesIO(pptx_bytes))

        slides_with_text = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slides_with_text += 1
                        break

        assert slides_with_text >= len(prs.slides) * 0.8, (
            f"至少 80% 投影片應有文字，實際: {slides_with_text}/{len(prs.slides)}"
        )
