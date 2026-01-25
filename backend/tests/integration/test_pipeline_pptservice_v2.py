"""
PPTServiceV2 Pipeline Integration Test

測試完整的 PPT 生成流程，透過 PPTServiceV2 執行所有階段。
"""

import io
import json
from datetime import datetime
from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from app.services.llm_service import get_llm_service
from app.services.ppt_service_v2 import PPTServiceV2

TEST_USER_INPUT = """
# 機器學習簡介

機器學習是人工智慧的一個子領域，讓電腦能從資料中學習並做出預測。

## 主要類型
1. 監督式學習 - 從標記資料學習
2. 非監督式學習 - 發現資料中的模式
3. 強化學習 - 透過獎勵機制學習

## 應用領域
- 圖像辨識
- 自然語言處理
- 推薦系統
"""


@pytest.fixture
def llm_service():
    return get_llm_service()


@pytest.fixture
def ppt_service(llm_service):
    templates_path = Path(__file__).parent.parent.parent / "data" / "templates"
    return PPTServiceV2(llm_service=llm_service, templates_path=templates_path)


@pytest.fixture
def output_dir():
    output_path = Path(__file__).parent.parent / "outputs"
    output_path.mkdir(parents=True, exist_ok=True)
    return output_path


class TestPPTServiceV2Generate:
    """測試 PPTServiceV2.generate()"""

    @pytest.mark.asyncio
    async def test_generate_returns_valid_pptx_bytes(self, ppt_service: PPTServiceV2):
        """測試 generate() 回傳有效的 PPTX bytes"""
        result = await ppt_service.generate(
            user_input=TEST_USER_INPUT,
            template="modern_clean.pptx",
            slide_count=5,
            language="zh-TW",
        )

        assert isinstance(result, bytes), "結果應為 bytes"
        assert len(result) > 0, "PPTX bytes 不應為空"

        prs = PptxPresentation(io.BytesIO(result))
        assert len(prs.slides) > 0, "PPTX 應包含投影片"

    @pytest.mark.asyncio
    async def test_generate_respects_slide_count(self, ppt_service: PPTServiceV2):
        """測試 generate() 遵守 slide_count"""
        target_count = 5

        result = await ppt_service.generate(
            user_input=TEST_USER_INPUT,
            slide_count=target_count,
            language="zh-TW",
        )

        prs = PptxPresentation(io.BytesIO(result))
        actual_count = len(prs.slides)

        assert abs(actual_count - target_count) <= 2, (
            f"投影片數量應接近 {target_count}，實際: {actual_count}"
        )

    @pytest.mark.asyncio
    async def test_generate_slides_have_content(self, ppt_service: PPTServiceV2):
        """測試生成的投影片有內容"""
        result = await ppt_service.generate(
            user_input=TEST_USER_INPUT,
            slide_count=5,
            language="zh-TW",
        )

        prs = PptxPresentation(io.BytesIO(result))

        slides_with_text = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slides_with_text += 1
                    break

        min_expected = len(prs.slides) * 0.8
        assert slides_with_text >= min_expected, (
            f"至少 {min_expected} 張投影片應有內容，實際: {slides_with_text}"
        )


class TestPPTServiceV2GenerateStream:
    """測試 PPTServiceV2.generate_stream()"""

    @pytest.mark.asyncio
    async def test_stream_yields_progress_updates(self, ppt_service: PPTServiceV2):
        """測試 generate_stream() 產生進度更新"""
        stages_seen = []

        async for update in ppt_service.generate_stream(
            user_input=TEST_USER_INPUT,
            slide_count=5,
            language="zh-TW",
        ):
            stages_seen.append(update.get("stage"))

        expected = [
            "template_analysis",
            "content_generation",
            "content_organization",
            "pptx_building",
            "completed",
        ]
        for stage in expected:
            assert stage in stages_seen, f"缺少階段: {stage}"

    @pytest.mark.asyncio
    async def test_stream_final_update_contains_result(self, ppt_service: PPTServiceV2):
        """測試最後一則更新包含結果"""
        final_update = None

        async for update in ppt_service.generate_stream(
            user_input=TEST_USER_INPUT,
            slide_count=5,
            language="zh-TW",
        ):
            final_update = update

        assert final_update.get("stage") == "completed"
        assert "result" in final_update
        assert "stats" in final_update

        prs = PptxPresentation(io.BytesIO(final_update["result"]))
        assert len(prs.slides) > 0


class TestPPTServiceV2GenerateToFile:
    """測試 PPTServiceV2.generate_to_file()"""

    @pytest.mark.asyncio
    async def test_generate_to_file_creates_file(self, ppt_service: PPTServiceV2, output_dir: Path):
        """測試 generate_to_file() 建立檔案"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = output_dir / f"pipeline_test_{timestamp}.pptx"

        result_path = await ppt_service.generate_to_file(
            user_input=TEST_USER_INPUT,
            output_path=output_path,
            slide_count=5,
            language="zh-TW",
        )

        assert result_path.exists()
        assert result_path.stat().st_size > 0

        prs = PptxPresentation(result_path)
        assert len(prs.slides) > 0


class TestPipelineDataFlow:
    """測試 Pipeline 資料流"""

    @pytest.mark.asyncio
    async def test_pipeline_preserves_topic_content(self, ppt_service: PPTServiceV2):
        """測試 Pipeline 保留主題內容"""
        keywords = ["機器學習", "監督", "非監督", "強化學習"]

        result = await ppt_service.generate(
            user_input=TEST_USER_INPUT,
            slide_count=5,
            language="zh-TW",
        )

        prs = PptxPresentation(io.BytesIO(result))

        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text + " "

        found = [kw for kw in keywords if kw in all_text]
        assert len(found) >= len(keywords) // 2, f"應包含主題關鍵字，找到: {found}"


class TestPipelineErrorHandling:
    """測試錯誤處理"""

    @pytest.mark.asyncio
    async def test_invalid_template_fallback(self, ppt_service: PPTServiceV2):
        """測試無效 Template 時的 fallback"""
        result = await ppt_service.generate(
            user_input=TEST_USER_INPUT,
            template="nonexistent.pptx",
            slide_count=5,
            language="zh-TW",
        )

        assert isinstance(result, bytes)
        prs = PptxPresentation(io.BytesIO(result))
        assert len(prs.slides) > 0
