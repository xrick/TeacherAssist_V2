"""
ML.md 整合測試

完整測試五階段 PPTX 生成流程：
1. TemplateAnalyzer - 分析 Template 結構
2. ContentGenerator - LLM 擴展使用者輸入
3. ContentOrganizerV2 - 組織內容到 Template 結構
4. ImageEnricher - 圖片搜尋與注入 (Pexels API)
5. SlideBuilder - 建構最終 PPTX (含 Layout Engine)

輸出: backend/output/ml_integration_test_YYYYMMDD_HHMMSS.pptx
"""

import asyncio
import json
import sys
from datetime import datetime
from pathlib import Path

# Add backend to path for imports
sys.path.insert(0, str(Path(__file__).parent.parent.parent))

from app.services.llm_service import get_llm_service
from app.services.ppt_service_v2 import PPTServiceV2


async def run_ml_integration_test():
    """執行 ML.md 整合測試"""

    print("=" * 60)
    print("ML.md 整合測試 - 五階段 PPTX 生成")
    print("=" * 60)

    # 讀取測試資料
    ml_path = Path(__file__).parent.parent / "data" / "ml.md"
    if not ml_path.exists():
        print(f"❌ 找不到測試資料: {ml_path}")
        return None

    user_input = ml_path.read_text(encoding="utf-8")
    print(f"✅ 讀取測試資料: {len(user_input)} 字元")

    # 設定輸出路徑
    output_dir = Path(__file__).parent.parent.parent / "output"
    output_dir.mkdir(parents=True, exist_ok=True)

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_pptx = output_dir / f"ml_integration_test_{timestamp}.pptx"
    report_json = output_dir / f"ml_integration_test_{timestamp}.json"

    # 建立服務
    llm_service = get_llm_service()
    templates_path = Path(__file__).parent.parent.parent / "data" / "templates"

    ppt_service = PPTServiceV2(
        llm_service=llm_service,
        templates_path=templates_path,
    )

    # 執行生成（帶進度追蹤）
    print("\n開始生成...")
    print("-" * 40)

    report = {
        "test_time": datetime.now().isoformat(),
        "input_file": str(ml_path),
        "input_length": len(user_input),
        "stages": [],
        "output": {},
    }

    start_time = datetime.now()
    final_result = None

    try:
        async for update in ppt_service.generate_stream(
            user_input=user_input,
            template="modern_clean.pptx",
            slide_count=4,  # 減少投影片數量以加快測試
            language="zh-TW",
            add_images=True,
            images_per_slide=1,
        ):
            stage = update.get("stage", "")
            progress = update.get("progress", 0)
            message = update.get("message", "")

            # 記錄階段
            report["stages"].append(
                {
                    "stage": stage,
                    "progress": progress,
                    "message": message,
                    "timestamp": datetime.now().isoformat(),
                }
            )

            # 顯示進度
            stage_icon = {
                "template_analysis": "🔍",
                "content_generation": "✨",
                "content_organization": "📋",
                "image_enrichment": "🖼️",
                "pptx_building": "🏗️",
                "completed": "✅",
                "error": "❌",
            }.get(stage, "📌")

            print(f"{stage_icon} [{progress:3.0f}%] {stage}: {message}")

            # 保存最終結果
            if "result" in update:
                final_result = update

    except Exception as e:
        print(f"❌ 生成失敗: {e}")
        report["error"] = str(e)
        report_json.write_text(json.dumps(report, ensure_ascii=False, indent=2))
        return None

    end_time = datetime.now()
    duration = (end_time - start_time).total_seconds()

    print("-" * 40)

    # 處理結果
    if final_result and "result" in final_result:
        pptx_bytes = final_result["result"]
        stats = final_result.get("stats", {})

        # 儲存 PPTX
        output_pptx.write_bytes(pptx_bytes)

        # 更新報告
        report["output"] = {
            "pptx_path": str(output_pptx),
            "pptx_size": len(pptx_bytes),
            "slide_count": stats.get("slide_count", 0),
            "duration_seconds": duration,
            "template": stats.get("template", ""),
        }

        # 分析 PPTX 內容
        try:
            from pptx import Presentation as PptxPresentation

            prs = PptxPresentation(output_pptx)

            slides_info = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                has_image = False

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            slide_text.append(text[:100])
                    if hasattr(shape, "image"):
                        has_image = True

                slides_info.append(
                    {
                        "index": i,
                        "text_preview": slide_text[:3],
                        "has_image": has_image,
                    }
                )

            report["output"]["slides_detail"] = slides_info
            report["output"]["slides_with_images"] = sum(1 for s in slides_info if s["has_image"])

        except Exception as e:
            print(f"⚠️ 分析 PPTX 失敗: {e}")

        # 儲存報告
        report_json.write_text(json.dumps(report, ensure_ascii=False, indent=2))

        print(f"\n✅ 測試完成!")
        print(f"   📊 投影片數量: {stats.get('slide_count', 'N/A')}")
        print(f"   ⏱️  執行時間: {duration:.2f} 秒")
        print(f"   📁 PPTX 大小: {len(pptx_bytes):,} bytes")
        print(f"\n📄 輸出檔案:")
        print(f"   PPTX: {output_pptx}")
        print(f"   報告: {report_json}")

        return output_pptx

    else:
        print("❌ 無生成結果")
        report_json.write_text(json.dumps(report, ensure_ascii=False, indent=2))
        return None


if __name__ == "__main__":
    result = asyncio.run(run_ml_integration_test())
    if result:
        print(f"\n🎉 成功產生: {result}")
    else:
        print("\n❌ 測試失敗")
        sys.exit(1)
